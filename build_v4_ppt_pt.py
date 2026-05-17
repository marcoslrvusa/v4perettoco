# -*- coding: utf-8 -*-
"""Generate V4 styled PPT (Portuguese) from unified markdown.
Uses python-pptx. Assumes markdown file at path.
"""
import re, sys
from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Paths
md_path = Path("/home/marcos/Desktop/v4perettoco-main-final/v4perettoco-main/squads/prime/clientes/atlas-copco-usa/campanhas/Estrategia-Integrada-Atlas-Copco-USA-UNIFICADO-PT.md")
output_path = Path("/home/marcos/Desktop/v4perettoco-main-final/v4perettoco-main/squads/prime/clientes/atlas-copco-usa/apresentacoes/AtlasCopco_Integrated_Strategy_V4_PT.pptx")

# Load markdown
text = md_path.read_text(encoding="utf-8")
# Simple parser: split sections by headings ##
sections = re.split(r"^##\s+", text, flags=re.MULTILINE)
# First element contains preamble, ignore
sections = sections[1:]

# Helper to create slide
def add_slide(prs, title_text, body_lines=None):
    blank_slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(blank_slide_layout)
    # background dark
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x11, 0x11, 0x11)
    # Title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(1)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = "Barlow"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xE3, 0x19, 0x19)
    p.alignment = PP_ALIGN.LEFT
    if body_lines:
        top_body = top + Inches(1.2)
        body_box = slide.shapes.add_textbox(left, top_body, width, Inches(5))
        tfb = body_box.text_frame
        for line in body_lines:
            p = tfb.add_paragraph()
            p.text = line
            p.font.name = "Barlow"
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
            p.alignment = PP_ALIGN.LEFT
    return slide

prs = Presentation()
# Title slide (from first heading line after #)
first_title_match = re.search(r"#\s+(.*)\[UNIFICADO\]", text)
if first_title_match:
    title = first_title_match.group(1).strip()
else:
    title = "Atlas Copco USA — Estratégia Integrada"
add_slide(prs, title)

# Map sections to slide titles
mapping = {
    "Sumário Executivo": "Sumário Executivo",
    "BUYER PERSONAS": "Personas",
    "DIAGNÓSTICO TÉCNICO": "Diagnóstico Técnico",
    "MAPA DE CALOR — SIMPLIFICADO & ACIONÁVEL": "Mapa de Calor",
    "PLANO DE AÇÃO — 90 DIAS": "Plano de Ação (90 dias)",
    "BUDGET & OKRs & ROI": "Budget & OKRs",
    "ESTRUTURA DE CAMPANHA": "Estrutura de Campanha",
    "MATRIZ DE DECISÃO": "Matriz de Decisão",
    "AVALIAÇÃO DE RISCOS": "Riscos",
}

for sec in sections:
    # first line is heading
    lines = sec.strip().splitlines()
    heading = lines[0].strip()
    body = lines[1:]
    # Clean markdown tables to simple lines
    cleaned = []
    for l in body:
        # replace markdown table rows with pipe separators trimmed
        if l.startswith('|'):
            # remove leading/trailing pipe and split
            cells = [c.strip() for c in l.strip('|').split('|')]
            cleaned.append(' | '.join(cells))
        else:
            cleaned.append(l)
    slide_title = mapping.get(heading.upper(), heading)
    add_slide(prs, slide_title, cleaned)

prs.save(output_path)
print(f"PPT generated at {output_path}")
